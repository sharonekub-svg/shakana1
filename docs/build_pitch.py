from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Brand colors
BG = RGBColor(0xF7, 0xF1, 0xE8)       # warm cream
CARD = RGBColor(0xFF, 0xFC, 0xF7)      # off-white
ACCENT = RGBColor(0xC9, 0x64, 0x42)    # terracotta/orange
DARK = RGBColor(0x2B, 0x21, 0x18)      # near-black
MID = RGBColor(0x6F, 0x62, 0x57)       # muted brown
LIGHT = RGBColor(0xE3, 0xD5, 0xC6)     # border beige
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2D, 0x9E, 0x6B)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, radius=False):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Arial"
    return txBox

def add_multiline(slide, lines, l, t, w, h, size=16, color=DARK, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for (text, bold, clr, sz) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(sz or size)
        run.font.bold = bold
        run.font.color.rgb = clr or color
        run.font.name = "Arial"
    return txBox

# ─────────────────────────────────────────────────────────────
# SLIDE 1 — TITLE
# ─────────────────────────────────────────────────────────────
slide = blank_slide(prs)
set_bg(slide, BG)
add_rect(slide, 0, 0, 13.33, 7.5, BG)

# big accent bar left
add_rect(slide, 0, 0, 0.18, 7.5, ACCENT)

# dot
add_rect(slide, 1.1, 2.2, 0.18, 0.18, ACCENT)

add_text(slide, "SHAKANA", 1.4, 1.6, 10, 1.0, size=52, bold=True, color=DARK)
add_text(slide, "קניות קבוצתיות לשכנים", 1.4, 2.65, 10, 0.7, size=24, color=MID)
add_text(slide, "Investor & Store Partner Deck  ·  2024", 1.4, 5.8, 10, 0.5, size=13, color=MID, italic=True)
add_text(slide, "shakana.app", 1.4, 6.5, 4, 0.5, size=13, color=ACCENT)

# ─────────────────────────────────────────────────────────────
# SLIDE 2 — WHO WE ARE
# ─────────────────────────────────────────────────────────────
slide = blank_slide(prs)
set_bg(slide, BG)
add_rect(slide, 0, 0, 0.18, 7.5, ACCENT)

add_text(slide, "WHO WE ARE", 0.5, 0.35, 5, 0.45, size=11, bold=True, color=ACCENT)
add_text(slide, "מי אנחנו", 0.5, 0.75, 10, 0.9, size=38, bold=True, color=DARK)

add_rect(slide, 0.5, 1.8, 12.3, 0.04, LIGHT)

add_text(slide, "Shakana is Israel's first social group-buying platform\nbuilt for apartment buildings and neighborhoods.", 0.5, 2.0, 12, 1.0, size=20, color=DARK)
add_text(slide, "We build the infrastructure for neighbors to shop together,\nsplit shipping costs, and unlock group discounts — from their phones.", 0.5, 3.1, 12, 1.0, size=18, color=MID)

# 3 pillars
for i, (title, sub) in enumerate([
    ("Social-first", "Built on WhatsApp sharing"),
    ("Mobile-first", "iOS + Android, Hebrew RTL"),
    ("Israel-first", "Designed for Israeli neighborhoods"),
]):
    x = 0.5 + i * 4.25
    add_rect(slide, x, 4.4, 3.9, 2.2, CARD)
    add_rect(slide, x, 4.4, 0.12, 2.2, ACCENT)
    add_text(slide, title, x + 0.3, 4.6, 3.4, 0.6, size=18, bold=True, color=DARK)
    add_text(slide, sub, x + 0.3, 5.2, 3.4, 0.9, size=14, color=MID)

# ─────────────────────────────────────────────────────────────
# SLIDE 3 — THE PROBLEM
# ─────────────────────────────────────────────────────────────
slide = blank_slide(prs)
set_bg(slide, BG)
add_rect(slide, 0, 0, 0.18, 7.5, ACCENT)

add_text(slide, "THE PROBLEM", 0.5, 0.35, 5, 0.45, size=11, bold=True, color=ACCENT)
add_text(slide, "הבעיה", 0.5, 0.75, 10, 0.9, size=38, bold=True, color=DARK)
add_rect(slide, 0.5, 1.8, 12.3, 0.04, LIGHT)

add_text(slide, "Israelis overpay for shipping — every single day.", 0.5, 2.0, 12, 0.7, size=20, bold=True, color=DARK)

for i, (num, label) in enumerate([
    ("₪35–60", "shipping per order, per person"),
    ("×12", "neighbors ordering the same store separately"),
    ("0", "products that connect them"),
]):
    x = 0.5 + i * 4.25
    add_rect(slide, x, 3.0, 3.9, 2.5, CARD)
    add_text(slide, num, x + 0.25, 3.2, 3.4, 0.9, size=34, bold=True, color=ACCENT)
    add_text(slide, label, x + 0.25, 4.1, 3.4, 0.8, size=14, color=MID)

add_text(slide, "Stores want volume. Customers want savings. Nobody connects them. Until now.", 0.5, 5.9, 12, 0.6, size=15, italic=True, color=MID)

# ─────────────────────────────────────────────────────────────
# SLIDE 4 — WHAT WE DO
# ─────────────────────────────────────────────────────────────
slide = blank_slide(prs)
set_bg(slide, BG)
add_rect(slide, 0, 0, 0.18, 7.5, ACCENT)

add_text(slide, "THE SOLUTION", 0.5, 0.35, 5, 0.45, size=11, bold=True, color=ACCENT)
add_text(slide, "מה אנחנו עושים", 0.5, 0.75, 10, 0.9, size=38, bold=True, color=DARK)
add_rect(slide, 0.5, 1.8, 12.3, 0.04, LIGHT)

steps = [
    ("1", "User finds a product"),
    ("2", "Opens a group → shares to building WhatsApp"),
    ("3", "Neighbors join, pick size/color/qty"),
    ("4", "Timer counts down"),
    ("5", "One order. One delivery. One address."),
    ("6", "Everyone saves. 🎉"),
]

for i, (num, text) in enumerate(steps):
    row = i // 3
    col = i % 3
    x = 0.5 + col * 4.25
    y = 2.1 + row * 2.3
    add_rect(slide, x, y, 3.9, 1.9, CARD)
    add_text(slide, num, x + 0.25, y + 0.15, 0.5, 0.55, size=22, bold=True, color=ACCENT)
    add_text(slide, text, x + 0.25, y + 0.7, 3.3, 0.9, size=14, color=DARK)

# ─────────────────────────────────────────────────────────────
# SLIDE 5 — WHY NOW
# ─────────────────────────────────────────────────────────────
slide = blank_slide(prs)
set_bg(slide, BG)
add_rect(slide, 0, 0, 0.18, 7.5, ACCENT)

add_text(slide, "MARKET OPPORTUNITY", 0.5, 0.35, 8, 0.45, size=11, bold=True, color=ACCENT)
add_text(slide, "למה עכשיו", 0.5, 0.75, 10, 0.9, size=38, bold=True, color=DARK)
add_rect(slide, 0.5, 1.8, 12.3, 0.04, LIGHT)

stats = [
    ("3.2M", "Apartments in Israel"),
    ("85%+", "WhatsApp penetration"),
    ("$50B+", "Pinduoduo (China) — same model"),
    ("$20B", "TikTok Shop GMV 2023"),
]
for i, (num, label) in enumerate(stats):
    x = 0.5 + (i % 2) * 6.4
    y = 2.0 + (i // 2) * 2.1
    add_rect(slide, x, y, 5.9, 1.7, CARD)
    add_text(slide, num, x + 0.25, y + 0.1, 5, 0.8, size=32, bold=True, color=ACCENT)
    add_text(slide, label, x + 0.25, y + 0.95, 5, 0.55, size=15, color=MID)

add_text(slide, "Israelis already coordinate purchases in building WhatsApp groups — manually.\nNobody has productized this. We are first.", 0.5, 6.35, 12, 0.8, size=14, italic=True, color=MID)

# ─────────────────────────────────────────────────────────────
# SLIDE 6 — PAYMENT MODEL
# ─────────────────────────────────────────────────────────────
slide = blank_slide(prs)
set_bg(slide, BG)
add_rect(slide, 0, 0, 0.18, 7.5, ACCENT)

add_text(slide, "PAYMENT MODEL", 0.5, 0.35, 8, 0.45, size=11, bold=True, color=ACCENT)
add_text(slide, "איך הפיימנטס עובד", 0.5, 0.75, 10, 0.9, size=38, bold=True, color=DARK)
add_rect(slide, 0.5, 1.8, 12.3, 0.04, LIGHT)

# Option 1
add_rect(slide, 0.5, 2.0, 5.9, 4.8, CARD)
add_rect(slide, 0.5, 2.0, 0.12, 4.8, ACCENT)
add_text(slide, "OPTION 1", 0.75, 2.1, 5, 0.4, size=10, bold=True, color=ACCENT)
add_text(slide, "Bit Flow", 0.75, 2.5, 5.2, 0.65, size=24, bold=True, color=DARK)
add_text(slide, "Zero friction. No credit card needed.", 0.75, 3.1, 5.2, 0.5, size=13, color=MID)

for line in [
    "• Each user sees their total in-app",
    "• App generates personal Bit payment link",
    "• Payment goes directly to store",
    "• All paid → order placed",
]:
    pass

add_multiline(slide, [
    ("• Each user sees their total in-app", False, MID, 13),
    ("• App generates personal Bit link", False, MID, 13),
    ("• Payment goes directly to store", False, MID, 13),
    ("• All confirmed → order placed", False, MID, 13),
    ("", False, MID, 8),
    ("✓ Israelis already use Bit", True, GREEN, 13),
    ("✓ No credit card required", True, GREEN, 13),
], 0.75, 3.6, 5.2, 2.8, size=13, color=MID)

# Option 2
add_rect(slide, 6.9, 2.0, 5.9, 4.8, CARD)
add_rect(slide, 6.9, 2.0, 0.12, 4.8, ACCENT)
add_text(slide, "OPTION 2  ★ RECOMMENDED", 7.15, 2.1, 5.4, 0.4, size=10, bold=True, color=ACCENT)
add_text(slide, "Prepay & Auto-Order", 7.15, 2.5, 5.2, 0.65, size=22, bold=True, color=DARK)
add_text(slide, "Fully automated. Scales to infinity.", 7.15, 3.1, 5.2, 0.5, size=13, color=MID)

add_multiline(slide, [
    ("• User joins → picks size/color/qty", False, MID, 13),
    ("• User pays immediately (card/Apple Pay)", False, MID, 13),
    ("• Timer counts down (24–48h)", False, MID, 13),
    ("• Timer hits zero → order auto-placed", False, MID, 13),
    ("• One delivery to founder's address", False, MID, 13),
    ("", False, MID, 8),
    ("✓ Zero human action needed", True, GREEN, 13),
    ("✓ Full refund if group fails", True, GREEN, 13),
], 7.15, 3.6, 5.2, 2.8, size=13, color=MID)

# ─────────────────────────────────────────────────────────────
# SLIDE 7 — FOR STORES
# ─────────────────────────────────────────────────────────────
slide = blank_slide(prs)
set_bg(slide, BG)
add_rect(slide, 0, 0, 0.18, 7.5, ACCENT)

add_text(slide, "FOR STORE PARTNERS", 0.5, 0.35, 8, 0.45, size=11, bold=True, color=ACCENT)
add_text(slide, "למה חנויות בוחרות בנו", 0.5, 0.75, 10, 0.9, size=38, bold=True, color=DARK)
add_rect(slide, 0.5, 1.8, 12.3, 0.04, LIGHT)

benefits = [
    ("Bulk orders", "One order per building — not 12 individual ones"),
    ("Zero marketing spend", "Users bring their own neighbors"),
    ("Paid upfront", "No payment chasing, no credit risk"),
    ("Zero returns on sizing", "Users chose their own specs"),
    ("Neighborhood analytics", "Know which products trend per area"),
    ("New acquisition channel", "Every group order = word-of-mouth event"),
]

for i, (title, desc) in enumerate(benefits):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.4
    y = 2.05 + row * 1.6
    add_rect(slide, x, y, 5.9, 1.35, CARD)
    add_text(slide, title, x + 0.25, y + 0.1, 5.3, 0.5, size=15, bold=True, color=DARK)
    add_text(slide, desc, x + 0.25, y + 0.6, 5.3, 0.55, size=13, color=MID)

# ─────────────────────────────────────────────────────────────
# SLIDE 8 — BUSINESS MODEL
# ─────────────────────────────────────────────────────────────
slide = blank_slide(prs)
set_bg(slide, BG)
add_rect(slide, 0, 0, 0.18, 7.5, ACCENT)

add_text(slide, "BUSINESS MODEL", 0.5, 0.35, 8, 0.45, size=11, bold=True, color=ACCENT)
add_text(slide, "מודל עסקי", 0.5, 0.75, 10, 0.9, size=38, bold=True, color=DARK)
add_rect(slide, 0.5, 1.8, 12.3, 0.04, LIGHT)

streams = [
    ("3–5%", "Platform commission per completed group order"),
    ("₪199/mo", "Store subscription — analytics + featured placement"),
    ("CPM", "Promoted product listings in neighborhood feeds"),
    ("Data", "Neighborhood demand reports sold to brands"),
]
for i, (num, desc) in enumerate(streams):
    x = 0.5 + (i % 2) * 6.4
    y = 2.1 + (i // 2) * 1.85
    add_rect(slide, x, y, 5.9, 1.6, CARD)
    add_text(slide, num, x + 0.25, y + 0.1, 2, 0.7, size=28, bold=True, color=ACCENT)
    add_text(slide, desc, x + 0.25, y + 0.85, 5.3, 0.6, size=14, color=MID)

add_rect(slide, 0.5, 5.85, 12.3, 1.4, CARD)
add_text(slide, "Unit Economics Example", 0.75, 5.9, 6, 0.45, size=12, bold=True, color=DARK)
add_text(slide, "Avg group order ₪1,200  ×  4% commission  =  ₪48 per order  ·  1,000 orders/mo  =  ₪48,000 MRR at zero marginal cost", 0.75, 6.35, 12, 0.55, size=13, color=MID)

# ─────────────────────────────────────────────────────────────
# SLIDE 9 — TRACTION
# ─────────────────────────────────────────────────────────────
slide = blank_slide(prs)
set_bg(slide, BG)
add_rect(slide, 0, 0, 0.18, 7.5, ACCENT)

add_text(slide, "TRACTION & ROADMAP", 0.5, 0.35, 8, 0.45, size=11, bold=True, color=ACCENT)
add_text(slide, "איפה אנחנו עכשיו", 0.5, 0.75, 10, 0.9, size=38, bold=True, color=DARK)
add_rect(slide, 0.5, 1.8, 12.3, 0.04, LIGHT)

add_rect(slide, 0.5, 2.0, 3.8, 4.8, CARD)
add_text(slide, "TODAY  ✓", 0.75, 2.15, 3.3, 0.45, size=12, bold=True, color=GREEN)
for j, t in enumerate(["App built & deployed", "iOS + Android live", "Stripe escrow live", "Hebrew / RTL", "Email OTP auth", "First stores in talks"]):
    add_text(slide, f"✓  {t}", 0.75, 2.65 + j * 0.62, 3.3, 0.55, size=13, color=DARK)

add_rect(slide, 4.75, 2.0, 3.8, 4.8, CARD)
add_text(slide, "6 MONTHS", 5.0, 2.15, 3.3, 0.45, size=12, bold=True, color=ACCENT)
for j, t in enumerate(["50 buildings active", "200 group orders/mo", "10 store partners", "WhatsApp invite cards", "Auto-order flow live"]):
    add_text(slide, f"→  {t}", 5.0, 2.65 + j * 0.62, 3.3, 0.55, size=13, color=DARK)

add_rect(slide, 9.0, 2.0, 3.8, 4.8, CARD)
add_text(slide, "12 MONTHS", 9.25, 2.15, 3.3, 0.45, size=12, bold=True, color=ACCENT)
for j, t in enumerate(["500 buildings", "5 cities", "₪500K ARR", "Building mgmt deals", "Series A ready"]):
    add_text(slide, f"→  {t}", 9.25, 2.65 + j * 0.62, 3.3, 0.55, size=13, color=DARK)

# ─────────────────────────────────────────────────────────────
# SLIDE 10 — THE ASK
# ─────────────────────────────────────────────────────────────
slide = blank_slide(prs)
set_bg(slide, BG)
add_rect(slide, 0, 0, 0.18, 7.5, ACCENT)

add_text(slide, "THE ASK", 0.5, 0.35, 8, 0.45, size=11, bold=True, color=ACCENT)
add_text(slide, "מה אנחנו צריכים", 0.5, 0.75, 10, 0.9, size=38, bold=True, color=DARK)
add_rect(slide, 0.5, 1.8, 12.3, 0.04, LIGHT)

add_rect(slide, 0.5, 2.1, 5.5, 2.2, CARD)
add_text(slide, "Raising", 0.75, 2.2, 4.5, 0.5, size=14, color=MID)
add_text(slide, "$300,000", 0.75, 2.65, 4.5, 0.9, size=40, bold=True, color=ACCENT)
add_text(slide, "Seed Round  ·  8% equity\nPre-money valuation: $3.25M", 0.75, 3.6, 4.5, 0.6, size=13, color=MID)

allocations = [
    ("40%  $120K", "Product & Engineering"),
    ("25%  $75K", "Store Partnerships & Sales"),
    ("20%  $60K", "Growth & Marketing"),
    ("15%  $45K", "Legal, Ops, Infrastructure"),
]
for i, (pct, label) in enumerate(allocations):
    y = 2.1 + i * 1.1
    add_rect(slide, 6.5, y, 6.3, 0.9, CARD)
    add_text(slide, pct, 6.75, y + 0.1, 2.5, 0.5, size=16, bold=True, color=ACCENT)
    add_text(slide, label, 9.4, y + 0.1, 3.1, 0.5, size=14, color=DARK)

add_text(slide, "18 months runway  ·  First 50 store partners  ·  1,000 active building groups", 0.5, 6.5, 12, 0.6, size=14, italic=True, color=MID)

# ─────────────────────────────────────────────────────────────
# SLIDE 11 — WHY INVEST
# ─────────────────────────────────────────────────────────────
slide = blank_slide(prs)
set_bg(slide, BG)
add_rect(slide, 0, 0, 0.18, 7.5, ACCENT)

add_text(slide, "WHY INVEST NOW", 0.5, 0.35, 8, 0.45, size=11, bold=True, color=ACCENT)
add_text(slide, "למה להשקיע עכשיו", 0.5, 0.75, 10, 0.9, size=38, bold=True, color=DARK)
add_rect(slide, 0.5, 1.8, 12.3, 0.04, LIGHT)

reasons = [
    ("Untapped market", "No competitor in Israel doing this"),
    ("Network effects", "Every building that joins locks in neighbors"),
    ("Viral by design", "Every group order is a WhatsApp share event"),
    ("Asset-light", "No inventory. No warehouse. No logistics."),
    ("Built already", "This is not a deck. The app exists. The tech works."),
    ("Founder-market fit", "Built by Israelis, for Israelis, in apartment buildings"),
]
for i, (title, desc) in enumerate(reasons):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.4
    y = 2.1 + row * 1.65
    add_rect(slide, x, y, 5.9, 1.4, CARD)
    add_rect(slide, x, y, 0.12, 1.4, ACCENT)
    add_text(slide, title, x + 0.3, y + 0.12, 5.3, 0.5, size=16, bold=True, color=DARK)
    add_text(slide, desc, x + 0.3, y + 0.65, 5.3, 0.55, size=13, color=MID)

# ─────────────────────────────────────────────────────────────
# SLIDE 12 — CLOSING
# ─────────────────────────────────────────────────────────────
slide = blank_slide(prs)
set_bg(slide, BG)
add_rect(slide, 0, 0, 0.18, 7.5, ACCENT)
add_rect(slide, 0.18, 0, 13.15, 7.5, BG)

add_text(slide, "SHAKANA", 1.0, 2.0, 11, 1.2, size=60, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_text(slide, "קניות קבוצתיות לשכנים", 1.0, 3.2, 11, 0.8, size=26, color=MID, align=PP_ALIGN.CENTER)
add_rect(slide, 4.5, 4.2, 4.3, 0.05, ACCENT)
add_text(slide, "sharon@shakana.app  ·  shakana.app", 1.0, 4.5, 11, 0.6, size=16, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, "Built in Israel. For Israel.", 1.0, 6.5, 11, 0.5, size=13, italic=True, color=MID, align=PP_ALIGN.CENTER)

prs.save("/home/user/shakana1/docs/Shakana-Pitch-Deck.pptx")
print("Done!")
