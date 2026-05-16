from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy, os

# ── BRAND COLORS ──────────────────────────────────────────────
BLACK  = RGBColor(0x0D, 0x0D, 0x0D)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
RED    = RGBColor(0xD4, 0x1C, 0x1C)
GRAY   = RGBColor(0xF2, 0xF2, 0xF2)
MID    = RGBColor(0x6B, 0x6B, 0x6B)
BORDER = RGBColor(0xD9, 0xD9, 0xD9)

FONT = "Arial"
LOGO_PATH = "/home/user/shakana1/apps/mobile/assets/icon.png"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=WHITE):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=GRAY, line_color=None, line_w=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_w or 1)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=16, bold=False, color=BLACK,
        align=PP_ALIGN.RIGHT, italic=False, rtl=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if rtl:
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        pPr.set(qn('a:rtl'), '1') if not hasattr(pPr, '_rtl_set') else None
        try:
            pPr.set('rtl', '1')
        except Exception:
            pass
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    run.font.name    = FONT
    return tb

def logo(slide, l, t, size=0.7):
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(l), Inches(t), Inches(size), Inches(size))

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════
s = blank(prs); bg(s, WHITE)
rect(s, 0, 0, 13.33, 0.08, fill=RED)       # top red stripe
rect(s, 0, 7.42, 13.33, 0.08, fill=BLACK)  # bottom stripe

logo(s, 0.6, 0.9, size=1.1)

txt(s, "שקאנה", 1.0, 1.0, 11.5, 1.8, size=72, bold=True, color=BLACK, align=PP_ALIGN.RIGHT)
txt(s, "קניות קבוצתיות לשכנים", 1.0, 2.8, 11.5, 0.8, size=28, color=MID, align=PP_ALIGN.RIGHT)

rect(s, 1.0, 3.9, 11.33, 0.04, fill=RED)

txt(s, "מצגת למשקיעים ולחנויות שותפות  ·  2024", 1.0, 4.1, 11.5, 0.6, size=15, color=MID, align=PP_ALIGN.RIGHT, italic=True)
txt(s, "shakana.app", 1.0, 6.6, 11.5, 0.5, size=14, color=RED, align=PP_ALIGN.RIGHT, bold=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — מי אנחנו
# ══════════════════════════════════════════════════════════════
s = blank(prs); bg(s, WHITE)
rect(s, 0, 0, 0.18, 7.5, fill=RED)
rect(s, 0, 0, 13.33, 0.05, fill=BLACK)

logo(s, 12.0, 0.25, size=0.7)

txt(s, "מי אנחנו", 0.5, 0.3, 12.0, 1.0, size=40, bold=True, color=BLACK, align=PP_ALIGN.RIGHT)
rect(s, 0.4, 1.45, 12.5, 0.04, fill=BORDER)

txt(s, "שקאנה היא הפלטפורמה הראשונה בישראל לקניות קבוצתיות עבור בניינים ושכונות.", 0.4, 1.6, 12.5, 0.8, size=20, color=BLACK, align=PP_ALIGN.RIGHT, bold=True)
txt(s, "אנחנו בונים את התשתית שמאפשרת לשכנים לקנות ביחד, לחלק עלויות משלוח, ולקבל הנחות קבוצתיות — ישירות מהטלפון.", 0.4, 2.4, 12.5, 0.9, size=17, color=MID, align=PP_ALIGN.RIGHT)

for i, (title, sub) in enumerate([
    ("חברתי קודם כל", "בנוי על שיתוף בוואטסאפ"),
    ("מובייל קודם כל", "iOS + Android, עברית ו-RTL"),
    ("ישראל קודם כל", "מתוכנן לבנייני מגורים ישראליים"),
]):
    x = 0.5 + i * 4.22
    rect(s, x, 3.5, 3.95, 2.8, fill=GRAY)
    rect(s, x, 3.5, 0.1, 2.8, fill=RED)
    txt(s, title, x + 0.25, 3.65, 3.5, 0.7, size=18, bold=True, color=BLACK, align=PP_ALIGN.RIGHT)
    txt(s, sub,   x + 0.25, 4.4,  3.5, 0.9, size=14, color=MID,   align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — הבעיה
# ══════════════════════════════════════════════════════════════
s = blank(prs); bg(s, WHITE)
rect(s, 0, 0, 0.18, 7.5, fill=RED)
rect(s, 0, 0, 13.33, 0.05, fill=BLACK)
logo(s, 12.0, 0.25, size=0.7)

txt(s, "הבעיה", 0.5, 0.3, 12.0, 1.0, size=40, bold=True, color=BLACK, align=PP_ALIGN.RIGHT)
rect(s, 0.4, 1.45, 12.5, 0.04, fill=BORDER)
txt(s, "ישראלים משלמים יותר מדי על משלוחים — כל יום.", 0.4, 1.6, 12.5, 0.7, size=20, bold=True, color=BLACK, align=PP_ALIGN.RIGHT)

for i, (num, label) in enumerate([
    ("₪35–60", "דמי משלוח לאדם, להזמנה"),
    ("×12",    "שכנים שמזמינים מאותה חנות בנפרד"),
    ("0",      "מוצרים שמחברים ביניהם"),
]):
    x = 0.5 + i * 4.22
    rect(s, x, 2.55, 3.95, 3.3, fill=GRAY)
    txt(s, num,   x + 0.25, 2.75, 3.5, 1.1, size=42, bold=True, color=RED,   align=PP_ALIGN.RIGHT)
    txt(s, label, x + 0.25, 3.9,  3.5, 1.2, size=15, color=BLACK, align=PP_ALIGN.RIGHT)

txt(s, "חנויות רוצות נפח. לקוחות רוצים חיסכון. אף אחד לא מחבר ביניהם. עד עכשיו.", 0.4, 6.1, 12.5, 0.7, size=14, color=MID, italic=True, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — מה אנחנו עושים
# ══════════════════════════════════════════════════════════════
s = blank(prs); bg(s, WHITE)
rect(s, 0, 0, 0.18, 7.5, fill=RED)
rect(s, 0, 0, 13.33, 0.05, fill=BLACK)
logo(s, 12.0, 0.25, size=0.7)

txt(s, "מה אנחנו עושים", 0.5, 0.3, 12.0, 1.0, size=40, bold=True, color=BLACK, align=PP_ALIGN.RIGHT)
rect(s, 0.4, 1.45, 12.5, 0.04, fill=BORDER)

steps = [
    ("1", "משתמש מוצא מוצר שהוא רוצה"),
    ("2", "פותח קבוצה → שולח לינק לקבוצת הבניין בוואטסאפ"),
    ("3", "שכנים מצטרפים ובוחרים מידה / צבע / כמות"),
    ("4", "טיימר סופר לאחור"),
    ("5", "הטיימר נגמר → הזמנה אחת נשלחת אוטומטית"),
    ("6", "משלוח אחד. חיסכון לכולם. 🎉"),
]
for i, (num, text) in enumerate(steps):
    col = i % 3; row = i // 3
    x = 0.4 + col * 4.27; y = 1.7 + row * 2.6
    rect(s, x, y, 4.0, 2.3, fill=GRAY)
    rect(s, x, y, 0.1, 2.3, fill=RED if i < 5 else BLACK)
    txt(s, num,  x + 0.25, y + 0.15, 3.5, 0.65, size=28, bold=True, color=RED, align=PP_ALIGN.RIGHT)
    txt(s, text, x + 0.25, y + 0.85, 3.5, 1.1,  size=14, color=BLACK, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — למה עכשיו
# ══════════════════════════════════════════════════════════════
s = blank(prs); bg(s, WHITE)
rect(s, 0, 0, 0.18, 7.5, fill=RED)
rect(s, 0, 0, 13.33, 0.05, fill=BLACK)
logo(s, 12.0, 0.25, size=0.7)

txt(s, "למה עכשיו", 0.5, 0.3, 12.0, 1.0, size=40, bold=True, color=BLACK, align=PP_ALIGN.RIGHT)
rect(s, 0.4, 1.45, 12.5, 0.04, fill=BORDER)

stats = [
    ("3.2M",  "דירות בישראל"),
    ("85%+",  "חדירת וואטסאפ בישראל"),
    ("$50B+", "Pinduoduo (סין) — אותו מודל"),
    ("$20B",  "TikTok Shop GMV 2023"),
]
for i, (num, label) in enumerate(stats):
    x = 0.4 + (i % 2) * 6.45; y = 1.7 + (i // 2) * 2.45
    rect(s, x, y, 6.1, 2.1, fill=GRAY)
    txt(s, num,   x + 0.25, y + 0.15, 5.5, 1.0, size=38, bold=True, color=RED,   align=PP_ALIGN.RIGHT)
    txt(s, label, x + 0.25, y + 1.2,  5.5, 0.65, size=16, color=BLACK, align=PP_ALIGN.RIGHT)

txt(s, "ישראלים כבר מתאמים קניות בקבוצות הוואטסאפ של הבניין — ידנית. אף אחד לא הפך את זה למוצר. אנחנו ראשונים.", 0.4, 6.3, 12.5, 0.8, size=14, color=MID, italic=True, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — מודל תשלום
# ══════════════════════════════════════════════════════════════
s = blank(prs); bg(s, WHITE)
rect(s, 0, 0, 0.18, 7.5, fill=RED)
rect(s, 0, 0, 13.33, 0.05, fill=BLACK)
logo(s, 12.0, 0.25, size=0.7)

txt(s, "מודל תשלום", 0.5, 0.3, 12.0, 1.0, size=40, bold=True, color=BLACK, align=PP_ALIGN.RIGHT)
rect(s, 0.4, 1.45, 12.5, 0.04, fill=BORDER)

# Option 1
rect(s, 0.4, 1.65, 5.95, 5.5, fill=GRAY)
rect(s, 0.4, 1.65, 0.12, 5.5, fill=BLACK)
txt(s, "אפשרות 1", 0.7, 1.75, 5.4, 0.45, size=11, bold=True, color=MID, align=PP_ALIGN.RIGHT)
txt(s, "ביט פלואו", 0.7, 2.15, 5.4, 0.7, size=26, bold=True, color=BLACK, align=PP_ALIGN.RIGHT)
txt(s, "אפס חיכוך. ללא כרטיס אשראי.", 0.7, 2.85, 5.4, 0.5, size=14, color=MID, align=PP_ALIGN.RIGHT, italic=True)
rect(s, 0.55, 3.4, 5.65, 0.03, fill=BORDER)
for j, line in enumerate([
    "• כל משתמש רואה את הסכום שלו",
    "• האפליקציה מייצרת לינק ביט אישי",
    "• התשלום עובר ישירות לחנות",
    "• כולם שילמו → ההזמנה יוצאת",
    "✓ ישראלים כבר משתמשים בביט",
    "✓ ללא כרטיס אשראי",
]):
    c = RED if "✓" in line else BLACK
    txt(s, line, 0.7, 3.55 + j * 0.46, 5.4, 0.42, size=13, color=c, align=PP_ALIGN.RIGHT,
        bold=("✓" in line))

# Option 2
rect(s, 6.9, 1.65, 5.95, 5.5, fill=BLACK)
rect(s, 6.9, 1.65, 0.12, 5.5, fill=RED)
txt(s, "אפשרות 2  ★ מומלצת", 7.15, 1.75, 5.4, 0.45, size=11, bold=True, color=RED, align=PP_ALIGN.RIGHT)
txt(s, "תשלום מראש + הזמנה אוטומטית", 7.15, 2.15, 5.4, 0.7, size=22, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
txt(s, "אוטומטי לחלוטין. גדל ללא הגבלה.", 7.15, 2.85, 5.4, 0.5, size=14, color=MID, align=PP_ALIGN.RIGHT, italic=True)
rect(s, 7.05, 3.4, 5.65, 0.03, fill=MID)
for j, line in enumerate([
    "• מצטרף → בוחר מידה / צבע / כמות",
    "• משלם מיד (כרטיס / Apple Pay)",
    "• טיימר סופר לאחור (24–48 שעות)",
    "• הטיימר נגמר → הזמנה יוצאת אוטומטית",
    "• משלוח אחד לכתובת יוזם הקבוצה",
    "✓ ללא פעולה אנושית",
    "✓ החזר מלא אם הקבוצה נכשלת",
]):
    c = RED if "✓" in line else WHITE
    txt(s, line, 7.15, 3.55 + j * 0.46, 5.4, 0.42, size=13, color=c, align=PP_ALIGN.RIGHT,
        bold=("✓" in line))

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — לחנויות
# ══════════════════════════════════════════════════════════════
s = blank(prs); bg(s, WHITE)
rect(s, 0, 0, 0.18, 7.5, fill=RED)
rect(s, 0, 0, 13.33, 0.05, fill=BLACK)
logo(s, 12.0, 0.25, size=0.7)

txt(s, "למה חנויות בוחרות בנו", 0.5, 0.3, 12.0, 1.0, size=40, bold=True, color=BLACK, align=PP_ALIGN.RIGHT)
rect(s, 0.4, 1.45, 12.5, 0.04, fill=BORDER)

benefits = [
    ("הזמנות בנפח",       "הזמנה אחת לבניין — לא 12 הזמנות בנפרד"),
    ("אפס הוצאות שיווק",  "המשתמשים מביאים את השכנים שלהם"),
    ("תשלום מראש",        "אין רדיפה אחרי תשלומים, אין סיכון אשראי"),
    ("אפס החזרות גודל",   "כל משתמש בחר את המפרט שלו בעצמו"),
    ("אנליטיקת שכונה",    "דעו אילו מוצרים טרנדיים בכל אזור"),
    ("ערוץ רכישה חדש",   "כל הזמנה קבוצתית = אירוע שיווק עצמי"),
]
for i, (title, desc) in enumerate(benefits):
    col = i % 2; row = i // 2
    x = 0.4 + col * 6.45; y = 1.7 + row * 1.75
    rect(s, x, y, 6.1, 1.5, fill=GRAY)
    rect(s, x, y, 0.1, 1.5, fill=RED)
    txt(s, title, x + 0.3, y + 0.12, 5.5, 0.55, size=16, bold=True, color=BLACK, align=PP_ALIGN.RIGHT)
    txt(s, desc,  x + 0.3, y + 0.7,  5.5, 0.65, size=13, color=MID,   align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — מודל עסקי
# ══════════════════════════════════════════════════════════════
s = blank(prs); bg(s, WHITE)
rect(s, 0, 0, 0.18, 7.5, fill=RED)
rect(s, 0, 0, 13.33, 0.05, fill=BLACK)
logo(s, 12.0, 0.25, size=0.7)

txt(s, "מודל עסקי", 0.5, 0.3, 12.0, 1.0, size=40, bold=True, color=BLACK, align=PP_ALIGN.RIGHT)
rect(s, 0.4, 1.45, 12.5, 0.04, fill=BORDER)

streams = [
    ("3–5%",     "עמלת פלטפורמה על כל הזמנה קבוצתית שהושלמה"),
    ("₪199/חודש","מנוי חנות — אנליטיקה + מיקום מועדף"),
    ("CPM",      "רשימות מוצרים ממומנות בפידים של השכונה"),
    ("דאטה",     "דוחות ביקוש שכונתיים למותגים"),
]
for i, (num, desc) in enumerate(streams):
    x = 0.4 + (i % 2) * 6.45; y = 1.7 + (i // 2) * 2.1
    rect(s, x, y, 6.1, 1.85, fill=GRAY)
    txt(s, num,  x + 0.25, y + 0.12, 5.5, 0.85, size=34, bold=True, color=RED,   align=PP_ALIGN.RIGHT)
    txt(s, desc, x + 0.25, y + 1.05, 5.5, 0.65, size=14, color=BLACK, align=PP_ALIGN.RIGHT)

rect(s, 0.4, 6.0, 12.5, 1.1, fill=GRAY)
rect(s, 0.4, 6.0, 0.1,  1.1, fill=BLACK)
txt(s, "דוגמת יחידת כלכלית: הזמנה ממוצעת ₪1,200  ×  4% עמלה  =  ₪48 להזמנה  ·  1,000 הזמנות/חודש  =  ₪48,000 MRR בעלות שולית אפסית",
    0.65, 6.3, 12.1, 0.65, size=13, color=MID, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — Traction + רודמאפ
# ══════════════════════════════════════════════════════════════
s = blank(prs); bg(s, WHITE)
rect(s, 0, 0, 0.18, 7.5, fill=RED)
rect(s, 0, 0, 13.33, 0.05, fill=BLACK)
logo(s, 12.0, 0.25, size=0.7)

txt(s, "איפה אנחנו ולאן אנחנו הולכים", 0.5, 0.3, 12.0, 1.0, size=36, bold=True, color=BLACK, align=PP_ALIGN.RIGHT)
rect(s, 0.4, 1.45, 12.5, 0.04, fill=BORDER)

cols = [
    ("עכשיו  ✓", BLACK, [
        "✓  האפליקציה בנויה ובפרודקשן",
        "✓  iOS + Android",
        "✓  Stripe escrow פעיל",
        "✓  עברית ו-RTL",
        "✓  OTP אימות אימייל",
        "✓  חנויות ראשונות בשיחות",
    ]),
    ("6 חודשים", RED, [
        "→  50 בניינים פעילים",
        "→  200 הזמנות קבוצתיות/חודש",
        "→  10 חנויות שותפות",
        "→  כרטיסי שיתוף וואטסאפ",
        "→  הזמנה אוטומטית חי",
    ]),
    ("12 חודשים", MID, [
        "→  500 בניינים",
        "→  5 ערים",
        "→  ₪500K ARR",
        "→  הסכמי ניהול בניינים",
        "→  מוכנים לסיבוב A",
    ]),
]
for i, (title, color, items) in enumerate(cols):
    x = 0.4 + i * 4.27
    rect(s, x, 1.65, 4.0, 5.55, fill=GRAY)
    rect(s, x, 1.65, 0.1,  5.55, fill=color)
    txt(s, title, x + 0.25, 1.78, 3.55, 0.55, size=15, bold=True, color=color, align=PP_ALIGN.RIGHT)
    for j, item in enumerate(items):
        txt(s, item, x + 0.25, 2.45 + j * 0.72, 3.55, 0.62, size=13, color=BLACK, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# SLIDE 10 — מה אנחנו צריכים (THE ASK)
# ══════════════════════════════════════════════════════════════
s = blank(prs); bg(s, WHITE)
rect(s, 0, 0, 0.18, 7.5, fill=RED)
rect(s, 0, 0, 13.33, 0.05, fill=BLACK)
logo(s, 12.0, 0.25, size=0.7)

txt(s, "מה אנחנו מחפשים", 0.5, 0.3, 12.0, 1.0, size=40, bold=True, color=BLACK, align=PP_ALIGN.RIGHT)
rect(s, 0.4, 1.45, 12.5, 0.04, fill=BORDER)

rect(s, 0.4, 1.65, 5.5, 3.0, fill=BLACK)
txt(s, "מגייסים", 0.6, 1.8, 5.1, 0.55, size=15, color=MID, align=PP_ALIGN.RIGHT)
txt(s, "$300,000", 0.6, 2.3, 5.1, 1.1, size=46, bold=True, color=RED, align=PP_ALIGN.RIGHT)
txt(s, "סיבוב ראשוני (Seed)  ·  8% אחוזי הון\nשווי חברה לפני כסף: $3.25M", 0.6, 3.45, 5.1, 0.9, size=13, color=WHITE, align=PP_ALIGN.RIGHT)

alloc = [
    ("40%  $120K", "פיתוח מוצר והנדסה"),
    ("25%  $75K",  "שותפויות חנויות ומכירות"),
    ("20%  $60K",  "צמיחה ושיווק"),
    ("15%  $45K",  "משפטי, תפעול, תשתית"),
]
for i, (pct, label) in enumerate(alloc):
    y = 1.65 + i * 1.2
    rect(s, 6.4, y, 6.4, 1.05, fill=GRAY)
    rect(s, 6.4, y, 0.1, 1.05, fill=RED)
    txt(s, pct,   6.65, y + 0.1, 3.0, 0.55, size=18, bold=True, color=RED, align=PP_ALIGN.RIGHT)
    txt(s, label, 9.7,  y + 0.1, 2.9, 0.55, size=15, color=BLACK, align=PP_ALIGN.RIGHT)

txt(s, "18 חודשי ראנוויי  ·  50 חנויות שותפות ראשונות  ·  1,000 קבוצות בניינים פעילות",
    0.4, 6.5, 12.5, 0.65, size=14, italic=True, color=MID, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# SLIDE 11 — למה להשקיע
# ══════════════════════════════════════════════════════════════
s = blank(prs); bg(s, WHITE)
rect(s, 0, 0, 0.18, 7.5, fill=RED)
rect(s, 0, 0, 13.33, 0.05, fill=BLACK)
logo(s, 12.0, 0.25, size=0.7)

txt(s, "למה להשקיע עכשיו", 0.5, 0.3, 12.0, 1.0, size=40, bold=True, color=BLACK, align=PP_ALIGN.RIGHT)
rect(s, 0.4, 1.45, 12.5, 0.04, fill=BORDER)

reasons = [
    ("שוק בתול",           "אין מתחרים בישראל שעושים את זה"),
    ("אפקט רשת",           "כל בניין שמצטרף נועל את השכנים"),
    ("ויראלי בתכנון",      "כל הזמנה קבוצתית = שיתוף בוואטסאפ"),
    ("ללא נכסים",           "אין מלאי. אין מחסן. אין לוגיסטיקה."),
    ("כבר בנוי",            "זו לא מצגת. האפליקציה קיימת. הטכנולוגיה עובדת."),
    ("התאמת מייסד-שוק",    "בנוי על ידי ישראלים, לישראלים, בבניינים"),
]
for i, (title, desc) in enumerate(reasons):
    col = i % 2; row = i // 2
    x = 0.4 + col * 6.45; y = 1.65 + row * 1.85
    rect(s, x, y, 6.1, 1.6, fill=GRAY)
    rect(s, x, y, 0.1, 1.6, fill=RED)
    txt(s, title, x + 0.3, y + 0.12, 5.5, 0.6,  size=17, bold=True, color=BLACK, align=PP_ALIGN.RIGHT)
    txt(s, desc,  x + 0.3, y + 0.75, 5.5, 0.65, size=13, color=MID,   align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# SLIDE 12 — CLOSING
# ══════════════════════════════════════════════════════════════
s = blank(prs); bg(s, WHITE)
rect(s, 0, 0, 13.33, 0.08, fill=RED)
rect(s, 0, 7.42, 13.33, 0.08, fill=BLACK)

logo(s, 6.1, 1.2, size=1.2)
txt(s, "שקאנה", 1.0, 2.5, 11.0, 1.8, size=68, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
txt(s, "קניות קבוצתיות לשכנים", 1.0, 4.3, 11.0, 0.75, size=24, color=MID, align=PP_ALIGN.CENTER)
rect(s, 4.5, 5.2, 4.3, 0.06, fill=RED)
txt(s, "sharon@shakana.app  ·  shakana.app", 1.0, 5.4, 11.0, 0.6, size=16, bold=True, color=RED, align=PP_ALIGN.CENTER)
txt(s, "בנוי בישראל. לישראל.", 1.0, 6.5, 11.0, 0.5, size=13, italic=True, color=MID, align=PP_ALIGN.CENTER)

# ── SAVE ──────────────────────────────────────────────────────
out = "/home/user/shakana1/docs/Shakana-Pitch-Deck-HE.pptx"
prs.save(out)
print(f"Saved → {out}")
